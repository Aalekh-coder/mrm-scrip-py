import os
import re
import json
import time
import tempfile
import requests
from urllib.parse import urlparse, urljoin
from dotenv import load_dotenv
from bs4 import BeautifulSoup

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

load_dotenv()

SCRAPERAPI_KEY = os.getenv("SCRAPER_API_KEY")

SCRAPERAPI_BASE = "http://api.scraperapi.com"
SCRAPERAPI_SEARCH_ENDPOINT = "https://api.scraperapi.com/structured/google/search"

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
        "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
    )
}

DOC_EXTENSIONS = (".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".txt")
CONTACT_PAGE_HINTS = ["contact", "contact-us", "contactus", "about", "about-us", "aboutus"]

# --- Regex patterns -------------------------------------------------

EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}")

# Loosely matches international/US style phone numbers
PHONE_RE = re.compile(
    r"(\+?\d{1,3}[\s.-]?)?(\(?\d{2,4}\)?[\s.-]?){2,4}\d{3,4}"
)

# Common social profile / username patterns
SOCIAL_RE = {
    "twitter_x": re.compile(r"(?:twitter\.com|x\.com)/([A-Za-z0-9_]{2,30})"),
    "instagram": re.compile(r"instagram\.com/([A-Za-z0-9_.]{2,30})"),
    "facebook": re.compile(r"facebook\.com/([A-Za-z0-9_.]{2,50})"),
    "linkedin": re.compile(r"linkedin\.com/(?:in|company)/([A-Za-z0-9_-]{2,60})"),
    "github": re.compile(r"github\.com/([A-Za-z0-9_-]{2,39})"),
    "telegram": re.compile(r"t\.me/([A-Za-z0-9_]{3,32})"),
    "youtube": re.compile(r"youtube\.com/(?:channel|c|@)?/?([A-Za-z0-9_-]{2,60})"),
    "whatsapp": re.compile(r"(?:wa\.me|whatsapp\.com)/([A-Za-z0-9_+]{5,20})"),
}

ADDRESS_HINTS = re.compile(
    r"\b\d{1,5}\s+[A-Za-z0-9.\s]{3,40}"
    r"(?:Street|St\.?|Avenue|Ave\.?|Road|Rd\.?|Boulevard|Blvd\.?|Lane|Ln\.?|"
    r"Drive|Dr\.?|Court|Ct\.?|Way|Suite|Floor)\b[A-Za-z0-9,.\s]{0,60}",
    re.IGNORECASE,
)


# --- Search ----------------------------------------------------------

def google_search(query, num_results=10):
    """Use ScraperAPI's structured Google Search endpoint to get top organic results."""
    params = {
        "api_key": SCRAPERAPI_KEY,
        "query": query,
        "num": num_results,
    }
    response = requests.get(SCRAPERAPI_SEARCH_ENDPOINT, params=params, timeout=30)
    response.raise_for_status()
    data = response.json()

    results = []
    for i, item in enumerate(data.get("organic_results", []), start=1):
        link = item.get("link") or item.get("url")
        if not link:
            continue
        results.append({
            "rank": i,
            "title": item.get("title"),
            "url": link,
            "displayed_link": item.get("displayed_link"),
            "snippet": item.get("snippet"),
        })
    return results[:num_results]


def is_domain_query(query):
    """Heuristic: does the query look like a bare domain name (e.g. example.com)?"""
    q = query.strip().lower()
    q = re.sub(r"^https?://", "", q)
    q = q.split("/")[0]
    return bool(re.match(r"^[a-z0-9-]+(\.[a-z0-9-]+)+$", q))


# --- Page fetching (via ScraperAPI proxy) ------------------------------

def fetch_page(url, timeout=60, render=False, premium=False):
    """Fetch a page through ScraperAPI and return (html, text, error).
    Automatically retries with render+premium enabled if the first attempt
    comes back 403 (common for JS-heavy / bot-protected sites like Instagram)."""
    params = {"api_key": SCRAPERAPI_KEY, "url": url}
    if render:
        params["render"] = "true"
    if premium:
        params["premium"] = "true"

    try:
        resp = requests.get(SCRAPERAPI_BASE, params=params, timeout=timeout)
        if resp.status_code == 403 and not (render and premium):
            # retry once with render + premium turned on
            retry_params = {**params, "render": "true", "premium": "true"}
            resp = requests.get(SCRAPERAPI_BASE, params=retry_params, timeout=timeout)
        resp.raise_for_status()
    except requests.RequestException as e:
        return None, None, str(e)

    html = resp.text
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    text = soup.get_text(separator=" ", strip=True)
    return html, text, None


def extract_contacts(html, text):
    emails = sorted(set(EMAIL_RE.findall(text)))
    phones = sorted(set(
        m.group(0).strip()
        for m in PHONE_RE.finditer(text)
        if len(re.sub(r"\D", "", m.group(0))) >= 7
    ))

    socials = {}
    for platform, pattern in SOCIAL_RE.items():
        matches = sorted(set(pattern.findall(html or "")))
        if matches:
            socials[platform] = matches

    addresses = sorted(set(m.strip() for m in ADDRESS_HINTS.findall(text)))

    return {
        "emails": emails,
        "phones": phones,
        "addresses": addresses,
        "socials": socials,
    }


def extract_doc_links(html, base_url):
    if not html:
        return []
    soup = BeautifulSoup(html, "html.parser")
    links = set()
    for a in soup.find_all("a", href=True):
        href = a["href"]
        if href.lower().endswith(DOC_EXTENSIONS):
            links.add(urljoin(base_url, href))
    return sorted(links)


def find_relevant_internal_links(html, base_url, hints):
    """Find About/Contact style links on a page, restricted to the same domain."""
    if not html:
        return []
    soup = BeautifulSoup(html, "html.parser")
    base_netloc = urlparse(base_url).netloc
    found = set()
    for a in soup.find_all("a", href=True):
        full = urljoin(base_url, a["href"])
        parsed = urlparse(full)
        if parsed.netloc and parsed.netloc != base_netloc:
            continue
        if any(hint in parsed.path.lower() for hint in hints):
            found.add(full)
    return sorted(found)


# --- Document (PDF/DOCX) scraping --------------------------------------

def extract_text_from_pdf(path):
    if pdfplumber is None:
        return ""
    chunks = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    chunks.append(page_text)
    except Exception as e:
        print(f"  [!] Failed to parse PDF {path}: {e}")
    return "\n".join(chunks)


def extract_text_from_docx(path):
    if docx is None:
        return ""
    try:
        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    except Exception as e:
        print(f"  [!] Failed to parse DOCX {path}: {e}")
        return ""


def extract_text_from_pptx(path):
    if Presentation is None:
        return ""
    chunks = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line:
                            chunks.append(line)
                # Tables on slides
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                chunks.append(cell.text)
                # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    chunks.append(notes)
    except Exception as e:
        print(f"  [!] Failed to parse PPTX {path}: {e}")
    return "\n".join(chunks)


def extract_text_from_xlsx(path):
    if openpyxl is None:
        return ""
    chunks = []
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        chunks.append(str(cell))
    except Exception as e:
        print(f"  [!] Failed to parse XLSX {path}: {e}")
    return "\n".join(chunks)


def extract_text_from_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"  [!] Failed to read TXT {path}: {e}")
        return ""


def scrape_document(doc_url):
    """Download a PDF/DOC(X) via ScraperAPI and extract contact details from its text."""
    try:
        resp = requests.get(
            SCRAPERAPI_BASE,
            params={"api_key": SCRAPERAPI_KEY, "url": doc_url},
            timeout=90,
        )
        resp.raise_for_status()
    except requests.RequestException as e:
        return {"url": doc_url, "status": "failed", "error": str(e)}

    suffix = os.path.splitext(urlparse(doc_url).path)[1].lower() or ".pdf"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(resp.content)
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(tmp_path)
        elif suffix in (".doc", ".docx"):
            text = extract_text_from_docx(tmp_path)
        elif suffix in (".ppt", ".pptx"):
            text = extract_text_from_pptx(tmp_path)
        elif suffix in (".xls", ".xlsx"):
            text = extract_text_from_xlsx(tmp_path)
        elif suffix == ".txt":
            text = extract_text_from_txt(tmp_path)
        else:
            text = ""
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    if not text:
        return {"url": doc_url, "status": "empty"}

    contacts = extract_contacts(html=None, text=text)
    return {"url": doc_url, "status": "ok", **contacts}


# --- Per-URL scraping ---------------------------------------------------

def scrape_url(url, domain_mode=False):
    domain = urlparse(url).netloc
    html, text, error = fetch_page(url)

    if error:
        return {
            "url": url,
            "domain": domain,
            "status": "failed",
            "error": error,
        }

    contacts = extract_contacts(html, text)
    result = {
        "url": url,
        "domain": domain,
        "status": "ok",
        **contacts,
    }

    doc_links = extract_doc_links(html, url)
    if doc_links:
        result["documents"] = [scrape_document(link) for link in doc_links[:5]]

    if domain_mode:
        sub_links = find_relevant_internal_links(html, url, CONTACT_PAGE_HINTS)
        if sub_links:
            about_contact = []
            for link in sub_links[:4]:
                sub_html, sub_text, sub_error = fetch_page(link)
                if sub_error:
                    about_contact.append({"url": link, "status": "failed", "error": sub_error})
                    continue
                sub_contacts = extract_contacts(sub_html, sub_text)
                sub_entry = {"url": link, "status": "ok", **sub_contacts}
                sub_docs = extract_doc_links(sub_html, link)
                if sub_docs:
                    sub_entry["documents"] = [scrape_document(d) for d in sub_docs[:5]]
                about_contact.append(sub_entry)
            result["about_contact_pages"] = about_contact

    return result


# --- Pipeline ----------------------------------------------------------

def search_and_scrape(query, num_results=10, delay=1.5):
    domain_mode = is_domain_query(query)
    search_results = google_search(query, num_results=num_results)

    enriched = []
    for item in search_results:
        url = item.get("url")
        if not url:
            continue

        scraped = scrape_url(url, domain_mode=domain_mode)
        enriched.append({**item, "contacts": scraped})

        time.sleep(delay)  # be polite, avoid hammering sites / rate limits

    return enriched


if __name__ == "__main__":
    if not SCRAPERAPI_KEY:
        raise SystemExit("ERROR: Set SCRAPERAPI_KEY in your .env file or environment.")

    query = input("Search Query: ")
    n = input("Number of results (default 10): ").strip()
    num_results = int(n) if n else 10

    data = search_and_scrape(query, num_results=num_results)
    print(json.dumps(data, indent=2, ensure_ascii=False))

    with open("results.json", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print("\nSaved to results.json")