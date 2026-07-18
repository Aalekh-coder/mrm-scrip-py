# import os
# import re
# import io
# import json
# import time
# import requests
# from urllib.parse import urlparse
# from dotenv import load_dotenv

# # Document parsing libs
# import pdfplumber
# from docx import Document as DocxDocument
# from pptx import Presentation
# import openpyxl
# import csv as csv_module
# from striprtf.striprtf import rtf_to_text

# load_dotenv()

# SCRAPERAPI_KEY = os.getenv("SCRAPER_API_KEY")

# SCRAPERAPI_BASE = "https://api.scraperapi.com/"
# SCRAPERAPI_SEARCH_URL = "https://api.scraperapi.com/structured/google/search"

# HEADERS = {
#     "User-Agent": (
#         "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
#         "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
#     )
# }

# # File types to search for, one "site:domain filetype:X" query per entry
# FILE_TYPES = [
#     "pdf", "doc", "docx", "ppt", "pptx",
#     "xls", "xlsx", "csv", "txt", "rtf",
# ]




# # --- Search (ScraperAPI structured Google search) ------------------------

# def google_search(query, num_results=None):
#     params = {
#         "api_key": SCRAPERAPI_KEY,
#         "query": query,
#     }
#     if num_results is not None:
#         params["num"] = num_results

#     response = requests.get(
#         SCRAPERAPI_SEARCH_URL,
#         params=params,
#         headers=HEADERS,
#         timeout=60,
#     )
#     response.raise_for_status()
#     data = response.json()

#     results = []
#     items = data.get("organic_results", [])
#     if num_results is not None:
#         items = items[:num_results]

#     for item in items:
#         link = item.get("link") or item.get("url")
#         if link:
#             results.append({
#                 "title": item.get("title"),
#                 "url": link,
#                 "snippet": item.get("snippet"),
#             })
#     return results


# def search_all_filetypes(domain, num_results_per_type=None, delay=1.5):
#     """Run one 'site:domain filetype:X' search per file type and collect all URLs."""
#     all_found = []
#     seen_urls = set()

#     for ext in FILE_TYPES:
#         query = f"site:{domain} filetype:{ext}"
#         print(f"Searching: {query}")
#         try:
#             results = google_search(query, num_results=num_results_per_type)
#         except Exception as e:
#             print(f"  [!] Search failed for {query}: {e}")
#             continue

#         for r in results:
#             if r["url"] not in seen_urls:
#                 seen_urls.add(r["url"])
#                 r["filetype"] = ext
#                 r["query"] = query
#                 all_found.append(r)

#         time.sleep(delay)

#     return all_found


# # --- Fetching through ScraperAPI proxy -----------------------------------

# def scraperapi_fetch(url, timeout=60):
#     params = {"api_key": SCRAPERAPI_KEY, "url": url}
#     try:
#         resp = requests.get(SCRAPERAPI_BASE, params=params, headers=HEADERS, timeout=timeout)
#         resp.raise_for_status()
#     except requests.RequestException as e:
#         return None, str(e)
#     return resp.content, None


# # --- Document text extraction -----------------------------------------

# def extract_text_from_pdf(content_bytes):
#     chunks = []
#     try:
#         with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
#             for page in pdf.pages:
#                 chunks.append(page.extract_text() or "")
#     except Exception as e:
#         return "", str(e)
#     return "\n".join(chunks), None


# def extract_text_from_docx(content_bytes):
#     try:
#         doc = DocxDocument(io.BytesIO(content_bytes))
#         parts = [p.text for p in doc.paragraphs]
#         for table in doc.tables:
#             for row in table.rows:
#                 for cell in row.cells:
#                     parts.append(cell.text)
#         return "\n".join(parts), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_pptx(content_bytes):
#     try:
#         prs = Presentation(io.BytesIO(content_bytes))
#         parts = []
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text") and shape.text:
#                     parts.append(shape.text)
#         return "\n".join(parts), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_xlsx(content_bytes):
#     try:
#         wb = openpyxl.load_workbook(io.BytesIO(content_bytes), data_only=True)
#         chunks = []
#         for ws in wb.worksheets:
#             for row in ws.iter_rows(values_only=True):
#                 for cell in row:
#                     if cell is not None:
#                         chunks.append(str(cell))
#         return "\n".join(chunks), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_csv(content_bytes):
#     try:
#         text = content_bytes.decode("utf-8", errors="ignore")
#         reader = csv_module.reader(io.StringIO(text))
#         chunks = [", ".join(row) for row in reader]
#         return "\n".join(chunks), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_txt(content_bytes):
#     try:
#         return content_bytes.decode("utf-8", errors="ignore"), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_rtf(content_bytes):
#     try:
#         raw = content_bytes.decode("utf-8", errors="ignore")
#         return rtf_to_text(raw), None
#     except Exception as e:
#         return "", str(e)


# def extract_text_from_document(ext, content_bytes):
#     ext = ext.lower()
#     if ext == "pdf":
#         return extract_text_from_pdf(content_bytes)
#     elif ext in ("doc", "docx"):
#         return extract_text_from_docx(content_bytes)
#     elif ext in ("ppt", "pptx"):
#         return extract_text_from_pptx(content_bytes)
#     elif ext in ("xls", "xlsx"):
#         return extract_text_from_xlsx(content_bytes)
#     elif ext == "csv":
#         return extract_text_from_csv(content_bytes)
#     elif ext == "txt":
#         return extract_text_from_txt(content_bytes)
#     elif ext == "rtf":
#         return extract_text_from_rtf(content_bytes)
#     return "", "unsupported document type"


# # --- Contact extraction ---------------------------------------------------


# def scrape_document(url, filetype):
#     content, error = scraperapi_fetch(url)
#     if error:
#         return {"url": url, "filetype": filetype, "status": "failed", "error": error}

#     text, doc_error = extract_text_from_document(filetype, content)
#     if doc_error and not text:
#         return {"url": url, "filetype": filetype, "status": "failed", "error": doc_error}

#     return {"url": url, "filetype": filetype, "status": "ok"}


# # --- Pipeline ----------------------------------------------------------

# def scrape_all_files_for_domain(domain, num_results_per_type=None, delay=1.5):
#     found = search_all_filetypes(domain, num_results_per_type=num_results_per_type, delay=delay)

#     print(f"\nFound {len(found)} unique document(s) across {len(FILE_TYPES)} file types. Scraping...\n")

#     results = []
#     for item in found:
#         print(f"Scraping: {item['url']}")
#         scraped = scrape_document(item["url"], item["filetype"])
#         results.append({
#             "title": item.get("title"),
#             "snippet": item.get("snippet"),
#             "query": item.get("query"),
#             **scraped,
#         })
#         time.sleep(delay)

#     return results


# if __name__ == "__main__":
#     if not SCRAPERAPI_KEY:
#         raise SystemExit("ERROR: Set SCRAPER_API_KEY in your .env file or environment.")

#     domain = input("Domain (e.g. example.com): ").strip()
#     n = input("Results per file type (leave blank for all): ").strip()
#     num_results_per_type = int(n) if n else None

#     data = scrape_all_files_for_domain(domain, num_results_per_type=num_results_per_type)
#     print(json.dumps(data, indent=2, ensure_ascii=False))

#     with open("results.json", "w", encoding="utf-8") as f:
#         json.dump(data, f, indent=2, ensure_ascii=False)
#     print(f"\nSaved {len(data)} document result(s) to results.json")



import os
import re
import io
import json
import time
import requests
from urllib.parse import urlparse
from dotenv import load_dotenv

# Document parsing libs
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import csv as csv_module
from striprtf.striprtf import rtf_to_text

load_dotenv()

SCRAPERAPI_KEY = os.getenv("SCRAPER_API_KEY")

SCRAPERAPI_BASE = "https://api.scraperapi.com/"
SCRAPERAPI_SEARCH_URL = "https://api.scraperapi.com/structured/google/search"

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
        "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
    )
}

# File types to search for, one "site:domain filetype:X" query per entry
FILE_TYPES = [
    "pdf", "doc", "docx", "ppt", "pptx",
    "xls", "xlsx", "csv", "txt", "rtf",
]


# --- Search (ScraperAPI structured Google search) ------------------------

def google_search(query, num_results=None):
    params = {
        "api_key": SCRAPERAPI_KEY,
        "query": query,
    }
    if num_results is not None:
        params["num"] = num_results

    response = requests.get(
        SCRAPERAPI_SEARCH_URL,
        params=params,
        headers=HEADERS,
        timeout=60,
    )
    response.raise_for_status()
    data = response.json()

    results = []
    items = data.get("organic_results", [])
    if num_results is not None:
        items = items[:num_results]

    for item in items:
        link = item.get("link") or item.get("url")
        if link:
            results.append({
                "title": item.get("title"),
                "url": link,
                "snippet": item.get("snippet"),
            })
    return results


def get_root_domain(domain):
    """Strip a subdomain down to its registrable root, e.g.
    admission.sgu.edu.in -> sgu.edu.in, www.example.com -> example.com.
    Simple heuristic (handles common multi-part TLDs like .edu.in/.co.in/.co.uk),
    not a full public suffix list, but good enough for typical cases."""
    parts = domain.lower().strip().split(".")
    if len(parts) <= 2:
        return domain
    multi_part_tlds = {"edu.in", "co.in", "co.uk", "org.in", "gov.in", "ac.in", "net.in"}
    last_two = ".".join(parts[-2:])
    if last_two in multi_part_tlds and len(parts) >= 3:
        return ".".join(parts[-3:])
    return last_two


def search_all_filetypes(domain, num_results_per_type=None, delay=1.5):
    """Run one 'site:domain filetype:X' search per file type and collect all URLs.
    If a subdomain search comes back empty, automatically retries against the
    root domain (e.g. admission.sgu.edu.in -> sgu.edu.in) since the files may
    be hosted on 'www.' or another subdomain instead."""
    all_found = []
    seen_urls = set()
    root_domain = get_root_domain(domain)
    fallback_needed = root_domain != domain

    for ext in FILE_TYPES:
        query = f"site:{domain} filetype:{ext}"
        print(f"Searching: {query}")
        try:
            results = google_search(query, num_results=num_results_per_type)
        except Exception as e:
            print(f"  [!] Search failed for {query}: {e}")
            results = []

        if not results and fallback_needed:
            fallback_query = f"site:{root_domain} filetype:{ext}"
            print(f"  -> no results for subdomain, retrying: {fallback_query}")
            try:
                results = google_search(fallback_query, num_results=num_results_per_type)
                query = fallback_query
            except Exception as e:
                print(f"  [!] Search failed for {fallback_query}: {e}")
                results = []

        for r in results:
            if r["url"] not in seen_urls:
                seen_urls.add(r["url"])
                r["filetype"] = ext
                r["query"] = query
                all_found.append(r)

        time.sleep(delay)

    return all_found


# --- Fetching through ScraperAPI proxy -----------------------------------

def scraperapi_fetch(url, timeout=60):
    params = {"api_key": SCRAPERAPI_KEY, "url": url}
    try:
        resp = requests.get(SCRAPERAPI_BASE, params=params, headers=HEADERS, timeout=timeout)
        resp.raise_for_status()
    except requests.RequestException as e:
        return None, str(e)
    return resp.content, None


# --- Document text extraction -----------------------------------------

def extract_text_from_pdf(content_bytes):
    chunks = []
    try:
        with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
            for page in pdf.pages:
                chunks.append(page.extract_text() or "")
    except Exception as e:
        return "", str(e)
    return "\n".join(chunks), None


def extract_text_from_docx(content_bytes):
    try:
        doc = DocxDocument(io.BytesIO(content_bytes))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        return "\n".join(parts), None
    except Exception as e:
        return "", str(e)


def extract_text_from_pptx(content_bytes):
    try:
        prs = Presentation(io.BytesIO(content_bytes))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts), None
    except Exception as e:
        return "", str(e)


def extract_text_from_xlsx(content_bytes):
    try:
        wb = openpyxl.load_workbook(io.BytesIO(content_bytes), data_only=True)
        chunks = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        chunks.append(str(cell))
        return "\n".join(chunks), None
    except Exception as e:
        return "", str(e)


def extract_text_from_csv(content_bytes):
    try:
        text = content_bytes.decode("utf-8", errors="ignore")
        reader = csv_module.reader(io.StringIO(text))
        chunks = [", ".join(row) for row in reader]
        return "\n".join(chunks), None
    except Exception as e:
        return "", str(e)


def extract_text_from_txt(content_bytes):
    try:
        return content_bytes.decode("utf-8", errors="ignore"), None
    except Exception as e:
        return "", str(e)


def extract_text_from_rtf(content_bytes):
    try:
        raw = content_bytes.decode("utf-8", errors="ignore")
        return rtf_to_text(raw), None
    except Exception as e:
        return "", str(e)


def extract_text_from_document(ext, content_bytes):
    ext = ext.lower()
    if ext == "pdf":
        return extract_text_from_pdf(content_bytes)
    elif ext in ("doc", "docx"):
        return extract_text_from_docx(content_bytes)
    elif ext in ("ppt", "pptx"):
        return extract_text_from_pptx(content_bytes)
    elif ext in ("xls", "xlsx"):
        return extract_text_from_xlsx(content_bytes)
    elif ext == "csv":
        return extract_text_from_csv(content_bytes)
    elif ext == "txt":
        return extract_text_from_txt(content_bytes)
    elif ext == "rtf":
        return extract_text_from_rtf(content_bytes)
    return "", "unsupported document type"


def scrape_document(url, filetype):
    content, error = scraperapi_fetch(url)
    if error:
        return {"url": url, "filetype": filetype, "status": "failed", "error": error}

    text, doc_error = extract_text_from_document(filetype, content)
    if doc_error and not text:
        return {"url": url, "filetype": filetype, "status": "failed", "error": doc_error}

    return {"url": url, "filetype": filetype, "status": "ok"}


# --- Pipeline ----------------------------------------------------------

def scrape_all_files_for_domain(domain, num_results_per_type=None, delay=1.5):
    found = search_all_filetypes(domain, num_results_per_type=num_results_per_type, delay=delay)

    print(f"\nFound {len(found)} unique document(s) across {len(FILE_TYPES)} file types. Scraping...\n")

    results = []
    for item in found:
        print(f"Scraping: {item['url']}")
        scraped = scrape_document(item["url"], item["filetype"])
        results.append({
            "title": item.get("title"),
            "snippet": item.get("snippet"),
            "query": item.get("query"),
            **scraped,
        })
        time.sleep(delay)

    return results


if __name__ == "__main__":
    if not SCRAPERAPI_KEY:
        raise SystemExit("ERROR: Set SCRAPER_API_KEY in your .env file or environment.")

    domain = input("Domain (e.g. example.com): ").strip()
    n = input("Results per file type (leave blank for all): ").strip()
    num_results_per_type = int(n) if n else None

    data = scrape_all_files_for_domain(domain, num_results_per_type=num_results_per_type)
    print(json.dumps(data, indent=2, ensure_ascii=False))

    with open("results.json", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print(f"\nSaved {len(data)} document result(s) to results.json")